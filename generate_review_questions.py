import os
import requests
from pptx import Presentation

UPLOAD_DIR = "uploads"
OLLAMA_MODEL = "mistral"

def extract_slide_text(pptx_path):
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        if texts:
            slides.append(f"Slide {i+1}: " + " ".join(texts))
    return slides

def generate_review_questions(slide_texts, model=OLLAMA_MODEL, num_questions=10):
    content = (
        f"Generate {num_questions} open-ended review questions for students based on this PowerPoint content.\n\n"
        + "\n".join(slide_texts)
        + "\n\nOnly list the questions. No answers."
    )

    response = requests.post(
        "http://localhost:11434/api/generate",
        json={"model": model, "prompt": content, "stream": False}
    )

    if response.status_code == 200:
        return response.json()["response"]
    else:
        return f"❌ Failed to generate: {response.text}"

if __name__ == "__main__":
    pptx_files = [f for f in os.listdir(UPLOAD_DIR) if f.endswith(".pptx")]

    if not pptx_files:
        print("❌ No PowerPoint files found in the 'uploads' folder.")
    else:
        for pptx_file in pptx_files:
            print(f"\n📂 Processing {pptx_file}")
            path = os.path.join(UPLOAD_DIR, pptx_file)
            slides = extract_slide_text(path)
            if slides:
                result = generate_review_questions(slides[:15])  # Limit for efficiency
                print("\n📘 Generated Review Questions:\n")
                print(result)
            else:
                print("⚠️ No readable text found in slides.")
