import streamlit as st
import os
import shutil
import csv
import difflib
import requests
import pandas as pd
from pptx import Presentation
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from process_documents import load_documents, split_documents
from qa_chain import create_vectorstore, load_vectorstore, build_qa_chain

# === CONFIG ===
UPLOAD_DIR = "uploads"
VECTOR_PATH = "vectorstore"
CUSTOM_QA_FILE = "custom_qa.csv"
BLOCKED_QA_FILE = "blocked_quiz_questions.csv"
GENERATED_DIR = "generated_questions"

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(GENERATED_DIR, exist_ok=True)

# === Load Q&A ===
def load_custom_qna(csv_path=CUSTOM_QA_FILE):
    qna_dict = {}
    if os.path.exists(csv_path):
        with open(csv_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                if len(row) >= 2:
                    qna_dict[row[0].strip().lower()] = row[1].strip()
    return qna_dict

def find_best_match(user_input, qna_dict, threshold=0.85):
    user_input = user_input.strip().lower()
    best_match = difflib.get_close_matches(user_input, qna_dict.keys(), n=1, cutoff=threshold)
    return qna_dict[best_match[0]] if best_match else None

custom_qna = load_custom_qna()

def load_blocked_questions(csv_path=BLOCKED_QA_FILE):
    blocklist = set()
    if os.path.exists(csv_path):
        with open(csv_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            next(reader, None)
            for row in reader:
                if row and row[0].strip():
                    blocklist.add(row[0].strip().lower())
    return blocklist

def is_blocked_question(user_input, blocked_set, threshold=0.85):
    user_input = user_input.strip().lower()
    return bool(difflib.get_close_matches(user_input, blocked_set, n=1, cutoff=threshold))

blocked_questions = load_blocked_questions()

def extract_slide_text(pptx_path):
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        text = " ".join(
            para.text.strip()
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs if para.text.strip()
        )
        if text:
            slides.append(text)
    return slides

def generate_review_questions_ollama(slide_texts, num_questions=10):
    content = (
        f"Generate {num_questions} open-ended review questions for students based on this PowerPoint content.\n\n"
        + "\n".join(slide_texts[:15])
        + "\n\nOnly list the questions. No answers."
    )
    response = requests.post(
        "http://localhost:11434/api/generate",
        json={"model": "mistral", "prompt": content, "stream": False}
    )
    return response.json()["response"] if response.status_code == 200 else "❌ Failed to generate."

# === STREAMLIT SETUP ===
st.set_page_config(page_title="Instructor AI Assistant", layout="centered")
st.title("📘 Instructor AI Assistant")

# === RESET ===
if st.button("🔁 Reset All"):
    shutil.rmtree(UPLOAD_DIR, ignore_errors=True)
    shutil.rmtree(VECTOR_PATH, ignore_errors=True)
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    st.experimental_rerun()

# === DISPLAY UPLOADED FILES (DROPDOWN ONLY) ===
if os.listdir(UPLOAD_DIR):
    st.subheader("📁 Uploaded Files")
    selected_file = st.selectbox(
        "Currently loaded into the AI assistant:",
        sorted(os.listdir(UPLOAD_DIR)),
        key="uploaded_file_dropdown"
    )

# === LOAD MEMORY VECTORSTORE ===
vectorstore, qa = None, None
if os.path.exists(os.path.join(VECTOR_PATH, "index.faiss")):
    with st.spinner("🔁 Loading saved memory..."):
        vectorstore = load_vectorstore()
        qa = build_qa_chain(vectorstore)
elif os.listdir(UPLOAD_DIR):
    with st.spinner("📚 Processing files..."):
        docs = load_documents(UPLOAD_DIR)
        chunks = split_documents(docs)
        vectorstore = create_vectorstore(chunks)
        qa = build_qa_chain(vectorstore)
        st.success("✅ Assistant Ready")


# === ASK A QUESTION ===
if qa:
    query = st.text_input("What do you want to know?")
    if query:
        with st.spinner("🤖 Thinking..."):
            if is_blocked_question(query, blocked_questions):
                st.warning("❌ I'm not allowed to answer quiz/exam questions.")
            else:
                answer = find_best_match(query, custom_qna)
                st.write(answer if answer else qa.run(query))

# === REVIEW QUESTION GENERATOR ===
st.markdown("---")
st.header("📘 Generate Review Questions from PowerPoints")

pptx_files = [f for f in os.listdir(UPLOAD_DIR) if f.endswith(".pptx")]
if pptx_files:
    selected_pptx = st.selectbox("Select a PowerPoint file", pptx_files)
    if st.button("⚙️ Generate Review Questions"):
        with st.spinner("Generating..."):
            slides = extract_slide_text(os.path.join(UPLOAD_DIR, selected_pptx))
            result = generate_review_questions_ollama(slides)
            st.markdown("### 📋 Review Questions")
            for line in result.split("\n"):
                if line.strip(): st.write(f"- {line.strip()}")

# === QUIZ PREVIEW MODE ===
st.markdown("---")
st.header("📝 Quiz Preview Mode")

quiz_files = [f for f in os.listdir(UPLOAD_DIR) if f.endswith("_review.csv")]
if quiz_files:
    selected_quiz = st.selectbox("Choose a quiz file", quiz_files, key="quiz_select")
    quiz_path = os.path.join(UPLOAD_DIR, selected_quiz)

    try:
        df = pd.read_csv(quiz_path)

        required_cols = ['Question', 'Option A', 'Option B', 'Option C', 'Option D', 'Correct Option Index']
        if all(col in df.columns for col in required_cols):
            show_answers = st.checkbox("✅ Show Correct Answers")

            clean_df = df.dropna(subset=['Question', 'Option A', 'Option B', 'Option C', 'Option D'])

            for i, row in clean_df.iterrows():
                question = row['Question']
                # Clean phrasing like: "What does the slide mean by"
                if isinstance(question, str) and question.lower().startswith("what does the slide mean by"):
                    question = question.replace("What does the slide mean by", "", 1).strip(": ").capitalize()

                st.markdown(f"**Q{i+1}: {question}**")
                st.write(f"A) {row['Option A']}")
                st.write(f"B) {row['Option B']}")
                st.write(f"C) {row['Option C']}")
                st.write(f"D) {row['Option D']}")

                if show_answers:
                    try:
                        correct_letter = ['A', 'B', 'C', 'D'][int(row['Correct Option Index'])]
                        st.success(f"✔️ Correct Answer: {correct_letter}")
                    except Exception:
                        st.warning("⚠️ Invalid answer index for this question.")

                st.markdown("---")
        else:
            st.warning("⚠️ This file does not match the expected quiz format.")

    except Exception as e:
        st.error(f"❌ Error reading file: {e}")

# === DEBATE REBUTTAL ANALYSIS ===
st.markdown("---")
st.header("⚔️ Debate Rebuttal Analyzer")

debate_file = st.file_uploader("Upload a debate argument (PDF or DOCX)", type=["pdf", "docx"], key="debate_upload")

if debate_file:
    from langchain_community.document_loaders import UnstructuredPDFLoader, UnstructuredWordDocumentLoader
    file_path = os.path.join(UPLOAD_DIR, debate_file.name)
    with st.spinner("📤 Saving and preparing..."):
        with open(file_path, "wb") as f:
            f.write(debate_file.getbuffer())

    try:
        with st.spinner("📄 Extracting text..."):
            loader = UnstructuredPDFLoader(file_path) if file_path.endswith(".pdf") else UnstructuredWordDocumentLoader(file_path)
            pages = loader.load()
            full_text = "\n".join([p.page_content for p in pages])
    except Exception as e:
        st.error(f"❌ Could not process file: {e}")
        full_text = None

    if full_text and st.button("🧠 Analyze for Rebuttable Areas"):
        with st.spinner("Analyzing with Mistral..."):
            prompt = (
                "You are a debate coach. Identify 3–5 specific claims from this student-written argument that are vulnerable to rebuttal.\n"
                "- Quote the sentence.\n"
                "- Explain why it's weak (emotional, unsupported, vague, etc.).\n"
                "- DO NOT provide sources or citations.\n\n"
                f"{full_text}\n\n"
                "Respond as a list with quoted text and a critique."
            )
            response = requests.post("http://localhost:11434/api/generate", json={"model": "mistral", "prompt": prompt, "stream": False})
            if response.status_code == 200:
                feedback = response.json()["response"]
                disclaimer = (
                    "\n\n---\n\n"
                    "**Note:** You are required to find academic sources to support or refute the claims highlighted above. "
                    "Some of these claims may be inaccurate or oversimplified. Do not follow them blindly. This analysis is intended to generate ideas, "
                    "not provide definitive answers. You are still responsible for conducting your own research to verify the evidence."
                )
                full_output = feedback + disclaimer
                st.markdown(full_output)

                # Create PDF
                buffer = BytesIO()
                pdf = canvas.Canvas(buffer, pagesize=letter)
                text = pdf.beginText(40, 750)
                text.setFont("Helvetica", 11)
                for line in full_output.split("\n"):
                    if text.getY() < 40:
                        pdf.drawText(text)
                        pdf.showPage()
                        text = pdf.beginText(40, 750)
                        text.setFont("Helvetica", 11)
                    text.textLine(line.strip())
                pdf.drawText(text)
                pdf.save()
                buffer.seek(0)

                st.download_button("📥 Download Feedback as PDF", buffer, "debate_rebuttal_feedback.pdf", mime="application/pdf")
